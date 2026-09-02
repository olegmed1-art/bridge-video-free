#!/usr/bin/env python3
"""Reproducible structural, bridge-deal and optional render QA for PPTX artifacts."""
from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

RANKS = set("AKQJT98765432")
SUIT_SYMBOLS = {"♠": "S", "♥": "H", "♦": "D", "♣": "C"}
BOARD_TITLE_RE = re.compile(r"^СДАЧА\s+(\d+)\b")
SEAT_RE = re.compile(r"^([NESW])\s*·")


def _shape_text(shape):
    try:
        return (shape.text or "").strip()
    except Exception:
        return ""


def _font_sizes(shape):
    out = []
    if not getattr(shape, "has_text_frame", False):
        return out
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                out.append(run.font.size.pt)
    return out


def _approx_overflow_warning(shape):
    """Conservative warning only; never used as a hard pass/fail criterion."""
    if not getattr(shape, "has_text_frame", False):
        return None
    text = _shape_text(shape)
    if not text:
        return None
    sizes = _font_sizes(shape)
    if not sizes:
        return None
    font_pt = sum(sizes) / len(sizes)
    if font_pt <= 0:
        return None
    width_pt = shape.width / 12700
    height_pt = shape.height / 12700
    usable_w = max(width_pt - 8, 1)
    usable_h = max(height_pt - 6, 1)
    chars_per_line = max(int(usable_w / (font_pt * 0.52)), 1)
    line_capacity = max(int(usable_h / (font_pt * 1.18)), 1)
    explicit = text.splitlines() or [""]
    estimated_lines = sum(
        max(1, (len(line) + chars_per_line - 1) // chars_per_line)
        for line in explicit
    )
    if estimated_lines > line_capacity * 1.35:
        return {
            "estimated_lines": estimated_lines,
            "estimated_capacity": line_capacity,
            "font_pt": round(font_pt, 2),
            "text_preview": text[:120],
        }
    return None


def _extract_board(slide):
    texts = [_shape_text(shape) for shape in slide.shapes if _shape_text(shape)]
    board = None
    for text in texts:
        match = BOARD_TITLE_RE.match(text)
        if match:
            board = int(match.group(1))
            break
    if board is None:
        return None

    hands = {}
    for index, text in enumerate(texts):
        match = SEAT_RE.match(text)
        if not match:
            continue
        seat = match.group(1)
        suits = {}
        cursor = index + 1
        while cursor + 1 < len(texts) and len(suits) < 4:
            symbol = texts[cursor]
            if symbol in SUIT_SYMBOLS:
                ranks = texts[cursor + 1].strip()
                if ranks in {"-", "—", "–"}:
                    ranks = ""
                suits[SUIT_SYMBOLS[symbol]] = ranks
                cursor += 2
            else:
                if SEAT_RE.match(symbol):
                    break
                cursor += 1
        if len(suits) == 4:
            hands[seat] = ".".join(suits[suit] if suits[suit] else "-" for suit in "SHDC")
    return board, hands


def _validate_hand(hand):
    cards = []
    parts = hand.split(".")
    if len(parts) != 4:
        return [], "not_four_suits"
    for suit, ranks in zip("SHDC", parts):
        if ranks in {"", "-", "—", "–"}:
            continue
        for rank in ranks:
            if rank not in RANKS:
                return [], f"invalid_rank:{rank}"
            cards.append(suit + rank)
    if len(cards) != 13:
        return cards, f"card_count:{len(cards)}"
    return cards, None


def analyze(path: Path, min_font_pt=8.0, bridge_board_check=False, render=False):
    prs = Presentation(str(path))
    slide_width, slide_height = prs.slide_width, prs.slide_height
    errors = []
    warnings = []
    min_font = None
    off_canvas = []
    zero_size = []
    empty_slides = []

    for slide_no, slide in enumerate(prs.slides, 1):
        has_content = False
        for shape in slide.shapes:
            text = _shape_text(shape)
            if text or shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_content = True
            if shape.width <= 0 or shape.height <= 0:
                zero_size.append({"slide": slide_no, "shape_id": shape.shape_id})
            if (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > slide_width
                or shape.top + shape.height > slide_height
            ):
                off_canvas.append(
                    {"slide": slide_no, "shape_id": shape.shape_id, "name": shape.name}
                )
            font_sizes = _font_sizes(shape)
            if font_sizes:
                local_min = min(font_sizes)
                min_font = local_min if min_font is None else min(min_font, local_min)
                if local_min < min_font_pt:
                    warnings.append(
                        {
                            "type": "FONT_BELOW_THRESHOLD",
                            "slide": slide_no,
                            "shape_id": shape.shape_id,
                            "min_font_pt": round(local_min, 2),
                        }
                    )
            overflow = _approx_overflow_warning(shape)
            if overflow:
                warnings.append(
                    {
                        "type": "TEXT_OVERFLOW_RISK",
                        "slide": slide_no,
                        "shape_id": shape.shape_id,
                        **overflow,
                    }
                )
        if not has_content:
            empty_slides.append(slide_no)

    if off_canvas:
        errors.append(
            {"type": "OFF_CANVAS_SHAPES", "count": len(off_canvas), "items": off_canvas[:50]}
        )
    if zero_size:
        errors.append(
            {"type": "ZERO_SIZE_SHAPES", "count": len(zero_size), "items": zero_size[:50]}
        )
    if empty_slides:
        errors.append({"type": "EMPTY_SLIDES", "slides": empty_slides})

    board_summary = None
    if bridge_board_check:
        board_items = []
        board_numbers = []
        for slide_no, slide in enumerate(prs.slides, 1):
            extracted = _extract_board(slide)
            if not extracted:
                continue
            board, hands = extracted
            board_numbers.append(board)
            board_errors = []
            all_cards = []
            for seat in "NESW":
                if seat not in hands:
                    board_errors.append(f"missing_{seat}")
                    continue
                cards, error = _validate_hand(hands[seat])
                all_cards.extend(cards)
                if error:
                    board_errors.append(f"{seat}:{error}")
            if len(all_cards) != 52:
                board_errors.append(f"total_cards:{len(all_cards)}")
            if len(set(all_cards)) != len(all_cards):
                board_errors.append("duplicate_cards")
            board_items.append(
                {
                    "slide": slide_no,
                    "board": board,
                    "hands": hands,
                    "status": "PASS" if not board_errors else "FAIL",
                    "errors": board_errors,
                }
            )
        if board_numbers and len(set(board_numbers)) != len(board_numbers):
            errors.append({"type": "DUPLICATE_BOARD_NUMBERS", "boards": board_numbers})
        failed = [item for item in board_items if item["status"] == "FAIL"]
        if failed:
            errors.append(
                {"type": "BRIDGE_BOARD_VALIDATION_FAILED", "count": len(failed), "items": failed}
            )
        board_summary = {
            "board_slide_count": len(board_items),
            "board_numbers": board_numbers,
            "passed": sum(item["status"] == "PASS" for item in board_items),
            "failed": sum(item["status"] == "FAIL" for item in board_items),
        }

    render_summary = None
    if render:
        executable = shutil.which("libreoffice") or shutil.which("soffice")
        if not executable:
            warnings.append({"type": "RENDER_SKIPPED", "reason": "libreoffice_not_found"})
        else:
            try:
                import fitz

                with tempfile.TemporaryDirectory(prefix="pptx-qa-") as temp_dir:
                    completed = subprocess.run(
                        [
                            executable,
                            "--headless",
                            "--convert-to",
                            "pdf",
                            "--outdir",
                            temp_dir,
                            str(path),
                        ],
                        stdout=subprocess.PIPE,
                        stderr=subprocess.PIPE,
                        text=True,
                        timeout=180,
                    )
                    pdf = Path(temp_dir) / (path.stem + ".pdf")
                    if completed.returncode != 0 or not pdf.is_file():
                        errors.append(
                            {
                                "type": "RENDER_FAILED",
                                "returncode": completed.returncode,
                                "stdout": completed.stdout[-500:],
                                "stderr": completed.stderr[-500:],
                            }
                        )
                    else:
                        doc = fitz.open(pdf)
                        ratios = []
                        for page in doc:
                            pix = page.get_pixmap(matrix=fitz.Matrix(0.5, 0.5), alpha=False)
                            data = pix.samples
                            nonwhite = 0
                            for index in range(0, len(data), pix.n):
                                if min(data[index : index + 3]) < 245:
                                    nonwhite += 1
                            ratios.append(nonwhite / (pix.width * pix.height))
                        pages = len(doc)
                        doc.close()
                        blank_pages = [
                            index + 1 for index, ratio in enumerate(ratios) if ratio < 0.005
                        ]
                        if pages != len(prs.slides):
                            errors.append(
                                {
                                    "type": "RENDER_PAGE_COUNT_MISMATCH",
                                    "slides": len(prs.slides),
                                    "pages": pages,
                                }
                            )
                        if blank_pages:
                            errors.append({"type": "RENDER_BLANK_PAGES", "pages": blank_pages})
                        render_summary = {
                            "status": (
                                "PASS"
                                if not blank_pages and pages == len(prs.slides)
                                else "FAIL"
                            ),
                            "pages": pages,
                            "nonwhite_ratio_min": round(min(ratios), 6) if ratios else None,
                            "nonwhite_ratio_max": round(max(ratios), 6) if ratios else None,
                            "nonwhite_ratio_mean": (
                                round(sum(ratios) / len(ratios), 6) if ratios else None
                            ),
                            "blank_pages": blank_pages,
                        }
            except Exception as exc:
                errors.append({"type": "RENDER_EXCEPTION", "error": str(exc)})

    return {
        "schema": "bridge-pptx-artifact-qa-v1",
        "status": "PASS" if not errors else "FAIL",
        "file": path.name,
        "slide_count": len(prs.slides),
        "slide_size_emu": {"width": slide_width, "height": slide_height},
        "structural": {
            "off_canvas_shape_count": len(off_canvas),
            "zero_size_shape_count": len(zero_size),
            "empty_slide_count": len(empty_slides),
            "minimum_explicit_font_pt": round(min_font, 2) if min_font else None,
            "warning_count": len(warnings),
        },
        "bridge_board_check": board_summary,
        "render": render_summary,
        "errors": errors,
        "warnings": warnings,
    }


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--output")
    parser.add_argument("--min-font-pt", type=float, default=8)
    parser.add_argument("--bridge-board-check", action="store_true")
    parser.add_argument("--render", action="store_true")
    args = parser.parse_args()
    result = analyze(
        Path(args.pptx), args.min_font_pt, args.bridge_board_check, args.render
    )
    text = json.dumps(result, ensure_ascii=False, indent=2, sort_keys=True) + "\n"
    if args.output:
        Path(args.output).write_text(text, encoding="utf-8")
    else:
        print(text, end="")
    raise SystemExit(0 if result["status"] == "PASS" else 2)


if __name__ == "__main__":
    main()
