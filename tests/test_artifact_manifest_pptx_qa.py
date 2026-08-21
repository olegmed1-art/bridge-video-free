from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from tools.artifact_manifest_v1 import build_manifest, file_entry, verify_manifest
from tools.pptx_artifact_qa import analyze


def test_manifest_roundtrip_and_tamper_detection(tmp_path: Path):
    artifact = tmp_path / "sample.txt"
    artifact.write_text("bridge-school\n", encoding="utf-8")
    manifest = build_manifest(
        [file_entry("derived_text", artifact)],
        [],
        {"algorithm": "test-v1"},
    )
    assert verify_manifest(manifest)["status"] == "PASS"

    artifact.write_text("changed\n", encoding="utf-8")
    result = verify_manifest(manifest)
    assert result["status"] == "FAIL"
    assert {item["error"] for item in result["errors"]} >= {"SIZE_MISMATCH", "SHA256_MISMATCH"}


def _add_text(slide, text: str, left: float, top: float, width: float, height: float, font=12):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font)
    return box


def test_pptx_qa_accepts_one_complete_bridge_board(tmp_path: Path):
    path = tmp_path / "board.pptx"
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(slide, "СДАЧА 1 · TEST", 0.3, 0.2, 3.0, 0.4)
    hands = {
        "N": ["AKQJ", "T98", "765", "432"],
        "E": ["T987", "AKQ", "J98", "765"],
        "S": ["6543", "J76", "AKQ", "JT9"],
        "W": ["2", "5432", "T432", "AKQ8"],
    }
    x_positions = {"N": 5.0, "E": 9.0, "S": 5.0, "W": 1.0}
    y_positions = {"N": 0.8, "E": 2.5, "S": 4.2, "W": 2.5}
    symbols = ["♠", "♥", "♦", "♣"]
    for seat in "NESW":
        x, y = x_positions[seat], y_positions[seat]
        _add_text(slide, f"{seat} · 10 HCP", x, y, 1.3, 0.25)
        for idx, (symbol, ranks) in enumerate(zip(symbols, hands[seat])):
            _add_text(slide, symbol, x, y + 0.3 + idx * 0.28, 0.25, 0.2)
            _add_text(slide, ranks, x + 0.3, y + 0.3 + idx * 0.28, 0.8, 0.2)
    prs.save(path)

    result = analyze(path, bridge_board_check=True, render=False)
    assert result["status"] == "PASS"
    assert result["bridge_board_check"]["board_slide_count"] == 1
    assert result["bridge_board_check"]["passed"] == 1
    assert result["structural"]["off_canvas_shape_count"] == 0


def test_committed_tournament_manifest_verifies():
    manifest_path = Path("qa/manifests/tournament_30041_artifact_manifest_v1.json")
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    result = verify_manifest(manifest)
    assert result["status"] == "PASS"
    assert result["entry_count"] == 2
    assert result["local_files_checked"] == 1


def test_committed_tournament_qa_evidence_is_pass():
    report = json.loads(Path("qa/evidence/tournament_30041_round2_pptx_qa.json").read_text(encoding="utf-8"))
    assert report["status"] == "PASS"
    assert report["slide_count"] == 27
    assert report["bridge_board_check"] == {
        "board_numbers": list(range(1, 25)),
        "board_slide_count": 24,
        "failed": 0,
        "passed": 24,
    }
    assert report["render"]["status"] == "PASS"
    assert report["render"]["blank_pages"] == []
